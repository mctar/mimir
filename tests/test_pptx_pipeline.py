import asyncio
import os
import tempfile
import pytest
import sys

sys.path.insert(0, os.path.dirname(os.path.dirname(__file__)))

import db


@pytest.fixture
def tmp_db(tmp_path):
    path = str(tmp_path / "test.db")
    asyncio.run(db.init_db(path))
    yield path
    asyncio.run(db.close_db())


def test_recaps_has_pptx_columns(tmp_db):
    """Les 4 nouvelles colonnes existent après init_db."""
    import aiosqlite

    async def check():
        conn = await aiosqlite.connect(tmp_db)
        cursor = await conn.execute("PRAGMA table_info(recaps)")
        cols = {row[1] for row in await cursor.fetchall()}
        await conn.close()
        return cols

    cols = asyncio.run(check())
    assert "pptx_instructions" in cols
    assert "deck_spec_json" in cols
    assert "deck_spec_model" in cols
    assert "deck_spec_at" in cols


def test_get_pptx_data_returns_none_for_missing(tmp_db):
    """get_pptx_data retourne None si pas de recap."""
    result = asyncio.run(db.get_pptx_data("nonexistent-session"))
    assert result is None


def test_save_and_get_pptx_instructions(tmp_db):
    """Round-trip save/get des instructions."""
    sid = "test-session-001"

    async def setup():
        await db._db.execute(
            "INSERT INTO sessions (id, topic, created_at) VALUES (?, ?, ?)",
            (sid, "Test", 1000.0),
        )
        await db._db.execute(
            "INSERT INTO recaps (session_id, recap_json, model, created_at) VALUES (?, ?, ?, ?)",
            (sid, '{"elevator_pitch": "test"}', "test", 1000.0),
        )
        await db._db.commit()

    asyncio.run(setup())
    asyncio.run(db.save_pptx_instructions(sid, "Traduire en anglais"))
    result = asyncio.run(db.get_pptx_data(sid))
    assert result["instructions"] == "Traduire en anglais"
    assert result["deck_spec"] is None


def test_save_and_get_deck_spec(tmp_db):
    """Round-trip save/get du deck_spec."""
    sid = "test-session-002"
    spec = {"schema_version": 1, "slides": [{"layout": "cover", "slots": {"title": "Test"}}]}

    async def setup():
        await db._db.execute(
            "INSERT INTO sessions (id, topic, created_at) VALUES (?, ?, ?)",
            (sid, "Test", 1000.0),
        )
        await db._db.execute(
            "INSERT INTO recaps (session_id, recap_json, model, created_at) VALUES (?, ?, ?, ?)",
            (sid, '{"elevator_pitch": "test"}', "test", 1000.0),
        )
        await db._db.commit()

    asyncio.run(setup())
    asyncio.run(db.save_deck_spec(sid, spec, "gemini-2.5-flash"))
    result = asyncio.run(db.get_pptx_data(sid))
    assert result["deck_spec"]["schema_version"] == 1
    assert result["deck_spec"]["slides"][0]["layout"] == "cover"
    assert result["deck_spec_model"] == "gemini-2.5-flash"


def test_pptx_instructions_survive_recap_regeneration(tmp_db):
    """store_recap does not wipe pptx_instructions (UPSERT behaviour)."""
    sid = "test-session-003"

    async def setup():
        await db._db.execute(
            "INSERT INTO sessions (id, topic, created_at) VALUES (?, ?, ?)",
            (sid, "Test", 1000.0),
        )
        await db._db.commit()

    asyncio.run(setup())
    # First recap
    asyncio.run(db.store_recap(sid, {"elevator_pitch": "v1"}, "model-a"))
    # Save pptx instructions
    asyncio.run(db.save_pptx_instructions(sid, "Traduire en anglais"))
    # Regenerate recap (simulates user clicking "Regenerate Recap")
    asyncio.run(db.store_recap(sid, {"elevator_pitch": "v2"}, "model-b"))
    # pptx_instructions must survive
    result = asyncio.run(db.get_pptx_data(sid))
    assert result["instructions"] == "Traduire en anglais"


def test_assemble_pptx_creates_file(tmp_path):
    """_assemble_pptx génère un fichier .pptx sans lever d'exception."""
    import export as exp

    deck_spec = {
        "schema_version": 1,
        "slides": [
            {"layout": "cover",      "slots": {"title": "Test Session", "date": "23 April 2026", "duration": "45m"}},
            {"layout": "quote-large","slots": {"title": "Pitch", "body": "Une phrase de pitch."}},
            {"layout": "bullets",    "slots": {"title": "Points clés", "bullets": ["Point 1", "Point 2", "Point 3"]}},
            {"layout": "three-columns", "slots": {"title": "Connexions", "col1": "A ↔ B", "col2": "C ↔ D", "col3": "E ↔ F"}},
            {"layout": "concepts",   "slots": {"title": "Concepts", "terms": ["A", "B", "C"], "edges": ["A → B", "B → C"]}},
        ],
    }
    output = str(tmp_path / "test_out.pptx")
    exp._assemble_pptx(deck_spec, output)
    assert os.path.exists(output)
    assert os.path.getsize(output) > 1000


def test_assemble_pptx_unknown_layout_fallback(tmp_path):
    """Un layout inconnu ne lève pas d'exception (fallback bullets)."""
    import export as exp

    deck_spec = {
        "schema_version": 1,
        "slides": [
            {"layout": "invented-layout", "slots": {"title": "Fallback", "bullets": ["item"]}},
        ],
    }
    output = str(tmp_path / "fallback.pptx")
    exp._assemble_pptx(deck_spec, output)
    assert os.path.exists(output)


def test_assemble_pptx_all_layouts(tmp_path):
    """text-large and two-columns layouts render without exception."""
    import export as exp

    deck_spec = {
        "schema_version": 1,
        "slides": [
            {"layout": "text-large",   "slots": {"title": "Text", "body": "Some long body text here."}},
            {"layout": "two-columns",  "slots": {"title": "Two cols", "left": "Left content", "right": "Right content"}},
        ],
    }
    output = str(tmp_path / "all_layouts.pptx")
    exp._assemble_pptx(deck_spec, output)
    assert os.path.exists(output)
    assert os.path.getsize(output) > 1000


def test_generate_deck_spec_parses_valid_json():
    """generate_deck_spec parse correctement un JSON LLM valide."""
    import export as exp
    import unittest.mock as mock

    valid_spec = {
        "schema_version": 1,
        "slides": [
            {"layout": "cover", "slots": {"title": "Test", "date": "2026-04-23", "duration": "30m"}},
            {"layout": "bullets", "slots": {"title": "Points", "bullets": ["A", "B"]}},
        ],
    }
    import json

    async def fake_call(tier, system, user):
        return json.dumps(valid_spec)

    with mock.patch.object(exp, "_llm_call_slides", fake_call):
        chain = [{"provider": "test", "model": "test-model"}]
        result = asyncio.run(
            exp.generate_deck_spec(
                transcript="Texte de test",
                recap={"elevator_pitch": "pitch"},
                instructions=None,
                current_deck_spec=None,
                chain=chain,
            )
        )

    assert result["schema_version"] == 1
    assert len(result["slides"]) == 2
    assert result["slides"][0]["layout"] == "cover"


def test_structural_qa_passes_valid_deck():
    from export import _structural_qa
    deck_spec = {"slides": [
        {"layout": "cover",        "slots": {"title": "The orchestrator model wins"}},
        {"layout": "bullets",      "slots": {"title": "The market is shifting now"}},
        {"layout": "cards-3",      "slots": {"title": "Three pillars of value"}},
        {"layout": "two-columns",  "slots": {"title": "Delivery unlocks margin"}},
        {"layout": "divider",      "slots": {"title": "VALUE PROPOSITION"}},
    ]}
    result = _structural_qa(deck_spec)
    assert result["passed"] is True
    assert result["issues"] == []


def test_structural_qa_fails_consecutive_layouts():
    from export import _structural_qa
    deck_spec = {"slides": [
        {"layout": "bullets", "slots": {"title": "A"}},
        {"layout": "bullets", "slots": {"title": "B"}},
        {"layout": "bullets", "slots": {"title": "C"}},
        {"layout": "bullets", "slots": {"title": "D"}},  # 4th consecutive
    ]}
    result = _structural_qa(deck_spec)
    assert result["passed"] is False
    assert any("consecutive" in issue for issue in result["issues"])


def test_structural_qa_fails_gerundive_titles():
    from export import _structural_qa
    deck_spec = {"slides": [
        {"layout": "bullets", "slots": {"title": "Positioning the Org"}},
        {"layout": "bullets", "slots": {"title": "Building the Model"}},
        {"layout": "bullets", "slots": {"title": "Delivering Value"}},  # 3rd gerundive
    ]}
    result = _structural_qa(deck_spec)
    assert result["passed"] is False
    assert any("gerundive" in issue for issue in result["issues"])


def test_structural_qa_fails_too_few_slides():
    from export import _structural_qa
    deck_spec = {"slides": [
        {"layout": "cover", "slots": {"title": "T"}},
        {"layout": "bullets", "slots": {"title": "B"}},
    ]}
    result = _structural_qa(deck_spec)
    assert result["passed"] is False
    assert any("count" in issue for issue in result["issues"])


def test_structural_qa_fails_too_many_slides():
    from export import _structural_qa
    slides = [{"layout": "bullets", "slots": {"title": f"Slide {i}"}} for i in range(16)]
    result = _structural_qa({"slides": slides})
    assert result["passed"] is False
    assert any("count" in issue for issue in result["issues"])


def test_generate_deck_spec_retries_on_invalid_json():
    """generate_deck_spec réessaie une fois si le JSON est invalide."""
    import export as exp
    import unittest.mock as mock
    import json

    valid_spec = {"schema_version": 1, "slides": [{"layout": "bullets", "slots": {"title": "T", "bullets": []}}]}
    call_count = {"n": 0}

    async def fake_call(tier, system, user):
        call_count["n"] += 1
        if call_count["n"] == 1:
            return "not valid json {{{"
        return json.dumps(valid_spec)

    with mock.patch.object(exp, "_llm_call_slides", fake_call):
        chain = [{"provider": "test", "model": "test-model"}]
        result = asyncio.run(
            exp.generate_deck_spec(
                transcript="Texte",
                recap={},
                instructions=None,
                current_deck_spec=None,
                chain=chain,
            )
        )

    assert call_count["n"] == 2
    assert result["slides"][0]["layout"] == "bullets"


def test_generate_deck_spec_multi_tier_fallback():
    """Si le tier 1 échoue (exception réseau), le tier 2 est utilisé."""
    import export as exp
    import unittest.mock as mock
    import json

    valid_spec = {"schema_version": 1, "slides": [{"layout": "cover", "slots": {"title": "Fallback", "date": "", "duration": ""}}]}
    call_count = {"n": 0}

    async def fake_call(tier, system, user):
        call_count["n"] += 1
        if tier["provider"] == "failing":
            raise RuntimeError("network error")
        return json.dumps(valid_spec)

    chain = [{"provider": "failing", "model": "bad"}, {"provider": "ok", "model": "good"}]
    with mock.patch.object(exp, "_llm_call_slides", fake_call):
        result = asyncio.run(
            exp.generate_deck_spec(
                transcript="Texte",
                recap={},
                instructions=None,
                current_deck_spec=None,
                chain=chain,
            )
        )

    assert result["slides"][0]["layout"] == "cover"
    assert call_count["n"] == 2  # tier 1 failed, tier 2 succeeded


def test_generate_deck_spec_truncates_slides_to_15():
    """Un deck de plus de 15 slides est tronqué à 15."""
    import export as exp
    import unittest.mock as mock
    import json

    slides = [{"layout": "bullets", "slots": {"title": f"Slide {i}", "bullets": []}} for i in range(20)]
    oversized_spec = {"schema_version": 1, "slides": slides}

    async def fake_call(tier, system, user):
        return json.dumps(oversized_spec)

    with mock.patch.object(exp, "_llm_call_slides", fake_call):
        result = asyncio.run(
            exp.generate_deck_spec(
                transcript="Texte",
                recap={},
                instructions=None,
                current_deck_spec=None,
                chain=[{"provider": "test", "model": "test"}],
            )
        )

    assert len(result["slides"]) == 15


def test_generate_deck_spec_strips_markdown_fence():
    """Un JSON enveloppé dans des fences markdown est correctement parsé."""
    import export as exp
    import unittest.mock as mock
    import json

    spec = {"schema_version": 1, "slides": [{"layout": "bullets", "slots": {"title": "Fenced", "bullets": ["a"]}}]}

    async def fake_call(tier, system, user):
        return f"```json\n{json.dumps(spec)}\n```"

    with mock.patch.object(exp, "_llm_call_slides", fake_call):
        result = asyncio.run(
            exp.generate_deck_spec(
                transcript="Texte",
                recap={},
                instructions=None,
                current_deck_spec=None,
                chain=[{"provider": "test", "model": "test"}],
            )
        )

    assert result["slides"][0]["layout"] == "bullets"
    assert result["slides"][0]["slots"]["title"] == "Fenced"


def test_assemble_pptx_cards3(tmp_path):
    """cards-3 génère une slide avec titre + 3 paires heading/content."""
    import export as exp
    from pptx import Presentation

    deck_spec = {
        "schema_version": 1,
        "slides": [
            {
                "layout": "cards-3",
                "slots": {
                    "title": "Trois catégories",
                    "cards": [
                        {"heading": "Cat A", "content": "Contenu de A"},
                        {"heading": "Cat B", "content": "Contenu de B"},
                        {"heading": "Cat C", "content": "Contenu de C"},
                    ],
                },
            }
        ],
    }
    output = str(tmp_path / "cards3.pptx")
    exp._assemble_pptx(deck_spec, output)
    assert os.path.exists(output)
    assert os.path.getsize(output) > 1000
    prs = Presentation(output)
    all_text = " ".join(
        shape.text_frame.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "Cat A" in all_text
    assert "Contenu de B" in all_text
    assert "Cat C" in all_text


def test_assemble_pptx_cards4(tmp_path):
    """cards-4 génère une slide avec titre + 4 paires heading/content."""
    import export as exp
    from pptx import Presentation

    deck_spec = {
        "schema_version": 1,
        "slides": [
            {
                "layout": "cards-4",
                "slots": {
                    "title": "Quatre thèmes",
                    "cards": [
                        {"heading": "Thème 1", "content": "Desc 1"},
                        {"heading": "Thème 2", "content": "Desc 2"},
                        {"heading": "Thème 3", "content": "Desc 3"},
                        {"heading": "Thème 4", "content": "Desc 4"},
                    ],
                },
            }
        ],
    }
    output = str(tmp_path / "cards4.pptx")
    exp._assemble_pptx(deck_spec, output)
    assert os.path.exists(output)
    assert os.path.getsize(output) > 1000
    prs = Presentation(output)
    all_text = " ".join(
        shape.text_frame.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "Thème 2" in all_text
    assert "Desc 4" in all_text


def test_assemble_pptx_cards5(tmp_path):
    """cards-5 génère une slide avec titre + 5 paires heading/content."""
    import export as exp
    from pptx import Presentation

    deck_spec = {
        "schema_version": 1,
        "slides": [
            {
                "layout": "cards-5",
                "slots": {
                    "title": "Cinq acteurs",
                    "cards": [
                        {"heading": "Acteur 1", "content": "Rôle 1"},
                        {"heading": "Acteur 2", "content": "Rôle 2"},
                        {"heading": "Acteur 3", "content": "Rôle 3"},
                        {"heading": "Acteur 4", "content": "Rôle 4"},
                        {"heading": "Acteur 5", "content": "Rôle 5"},
                    ],
                },
            }
        ],
    }
    output = str(tmp_path / "cards5.pptx")
    exp._assemble_pptx(deck_spec, output)
    assert os.path.exists(output)
    assert os.path.getsize(output) > 1000
    prs = Presentation(output)
    all_text = " ".join(
        shape.text_frame.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "Acteur 3" in all_text
    assert "Rôle 5" in all_text


def test_assemble_pptx_cards4_rounded(tmp_path):
    """cards-4-rounded génère une slide avec 4 paires heading/content (variante arrondie)."""
    import export as exp
    from pptx import Presentation

    deck_spec = {
        "schema_version": 1,
        "slides": [
            {
                "layout": "cards-4-rounded",
                "slots": {
                    "title": "Quatre étapes",
                    "cards": [
                        {"heading": "Étape 1", "content": "Description étape 1"},
                        {"heading": "Étape 2", "content": "Description étape 2"},
                        {"heading": "Étape 3", "content": "Description étape 3"},
                        {"heading": "Étape 4", "content": "Description étape 4"},
                    ],
                },
            }
        ],
    }
    output = str(tmp_path / "cards4r.pptx")
    exp._assemble_pptx(deck_spec, output)
    assert os.path.exists(output)
    assert os.path.getsize(output) > 1000
    prs = Presentation(output)
    all_text = " ".join(
        shape.text_frame.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "Étape 2" in all_text
    assert "Description étape 4" in all_text


def test_cards_no_shared_tags(tmp_path):
    """No tag file should be referenced by more than one slide (shared tags corrupt PPTX)."""
    import zipfile
    from collections import defaultdict
    from xml.etree import ElementTree as ET
    import export as exp

    deck_spec = {
        "schema_version": 1,
        "slides": [
            {"layout": "cards-3", "slots": {"title": "T1", "cards": [
                {"heading": "A", "content": "a"},
                {"heading": "B", "content": "b"},
                {"heading": "C", "content": "c"},
            ]}},
            {"layout": "cards-4", "slots": {"title": "T2", "cards": [
                {"heading": "A", "content": "a"},
                {"heading": "B", "content": "b"},
                {"heading": "C", "content": "c"},
                {"heading": "D", "content": "d"},
            ]}},
        ],
    }
    output = str(tmp_path / "no_shared_tags.pptx")
    exp._assemble_pptx(deck_spec, output)

    TAGS_RELTYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/tags"
    NS = "http://schemas.openxmlformats.org/package/2006/relationships"
    tag_owners = defaultdict(list)
    with zipfile.ZipFile(output) as z:
        for name in z.namelist():
            if "slides/_rels/" in name:
                tree = ET.fromstring(z.read(name))
                for rel in tree.findall(f"{{{NS}}}Relationship"):
                    if rel.get("Type") == TAGS_RELTYPE:
                        tag_owners[rel.get("Target")].append(name)

    for tag, owners in tag_owners.items():
        assert len(owners) == 1, f"Tag {tag} shared by {len(owners)} slides: {owners}"


def test_assemble_pptx_no_template_slides(tmp_path):
    """Le PPTX généré ne doit contenir QUE les slides du deck_spec, pas les slides template."""
    import export as exp
    from pptx import Presentation

    deck_spec = {
        "schema_version": 1,
        "slides": [
            {"layout": "cover", "slots": {"title": "Mon titre", "date": "2026-04-24"}},
            {"layout": "bullets", "slots": {"title": "Points clés", "bullets": ["Point A", "Point B"]}},
        ],
    }
    output = str(tmp_path / "no_template.pptx")
    exp._assemble_pptx(deck_spec, output)

    prs = Presentation(output)
    assert len(prs.slides) == 2, (
        f"Expected 2 slides (one per deck_spec entry), got {len(prs.slides)}. "
        "Template Lorem-ipsum slides are leaking into the output."
    )


def test_format_recap_v3():
    """_format_recap produces labelled sections for V3 structured fields."""
    from export import _format_recap

    recap = {
        "schema_version": 3,
        "transcript_stats": {"duration_minutes": 90},
        "positioning": {
            "what_to_sell": ["End-to-end reinvention"],
            "why_now": ["Agentic operations"],
            "why_well_positioned": ["Tri-pod"],
            "to_whom": ["Global 2000 CEO"],
        },
        "value_proposition": {
            "what_we_do": ["Design value engines"],
            "how_we_do_it": ["Connecting functions"],
            "how_we_get_paid": ["Value-based contracts"],
        },
        "positioning_statement": "For Global 2000...",
        "scope_boundaries_non_goals": ["Not a functional pitch", "Not one-off"],
    }
    result = _format_recap(recap)
    assert "=== POSITIONING ===" in result
    assert "What to sell?" in result
    assert "End-to-end reinvention" in result
    assert "=== VALUE PROPOSITION ===" in result
    assert "Design value engines" in result
    assert "=== POSITIONING STATEMENT ===" in result
    assert "For Global 2000" in result
    assert "=== SCOPE / BOUNDARIES / NON-GOALS ===" in result
    assert "Not a functional pitch" in result
    # Internal metadata must NOT appear as content sections
    assert "Schema Version" not in result
    assert "Transcript Stats" not in result


def test_format_recap_unknown_keys():
    """_format_recap appends non-V3 keys as generic sections."""
    from export import _format_recap

    recap = {
        "elevator_pitch": "A sharp pitch",
        "key_takeaways": [{"topics": ["A", "B"], "insight": "They are connected"}],
    }
    result = _format_recap(recap)
    assert "ELEVATOR PITCH" in result
    assert "A sharp pitch" in result
    assert "KEY TAKEAWAYS" in result
    assert "A ↔ B" in result
    assert "They are connected" in result


def test_build_user_prompt_instructions_first():
    """Instructions block appears before RÉCAP block in the prompt."""
    from export import _build_user_prompt

    prompt = _build_user_prompt(
        recap={"elevator_pitch": "test"},
        transcript="some text",
        instructions="Translate everything to English",
        current_deck_spec=None,
    )
    instructions_pos = prompt.find("INSTRUCTIONS")
    recap_pos = prompt.find("RÉCAP")
    assert instructions_pos >= 0, "INSTRUCTIONS block missing"
    assert recap_pos >= 0, "RÉCAP block missing"
    assert instructions_pos < recap_pos, "INSTRUCTIONS must come before RÉCAP"
    assert "TOUTES les slides" in prompt
    assert "Translate everything to English" in prompt


def test_build_user_prompt_no_instructions():
    """No INSTRUCTIONS block when instructions is None or empty."""
    from export import _build_user_prompt

    prompt = _build_user_prompt(
        recap={"elevator_pitch": "test"},
        transcript="some text",
        instructions=None,
        current_deck_spec=None,
    )
    assert "INSTRUCTIONS" not in prompt
    assert "RÉCAP" in prompt

    prompt2 = _build_user_prompt(
        recap={"elevator_pitch": "test"},
        transcript="some text",
        instructions="   ",  # whitespace only
        current_deck_spec=None,
    )
    assert "INSTRUCTIONS" not in prompt2


def test_assemble_divider_layout(tmp_path):
    """divider layout produces a slide with title and number filled."""
    import export as exp
    from pptx import Presentation

    deck_spec = {
        "schema_version": 1,
        "slides": [
            {"layout": "divider", "slots": {"title": "POSITIONING", "number": "01"}},
        ],
    }
    output = str(tmp_path / "divider.pptx")
    exp._assemble_pptx(deck_spec, output)
    assert os.path.exists(output)
    prs = Presentation(output)
    assert len(prs.slides) == 1
    slide = prs.slides[0]
    all_text = " ".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
    assert "POSITIONING" in all_text
    assert "01" in all_text


def test_build_user_prompt_session_context():
    """CONTEXTE block appears when session_topic or session_date is provided."""
    from export import _build_user_prompt

    prompt = _build_user_prompt(
        recap={"elevator_pitch": "test"},
        transcript="some text",
        instructions=None,
        current_deck_spec=None,
        session_topic="Intelligent Operations",
        session_date="24 April 2026",
    )
    assert "CONTEXTE" in prompt
    assert "Intelligent Operations" in prompt
    assert "24 April 2026" in prompt

    # No CONTEXTE block when both are empty
    prompt2 = _build_user_prompt(
        recap={"elevator_pitch": "test"},
        transcript="some text",
        instructions=None,
        current_deck_spec=None,
    )
    assert "CONTEXTE" not in prompt2


def test_structural_qa_handles_null_title():
    from export import _structural_qa
    deck_spec = {"slides": [
        {"layout": "cover",   "slots": {"title": None}},
        {"layout": "bullets", "slots": {"title": "The orchestrator model wins"}},
        {"layout": "cards-3", "slots": {"title": "Three value dimensions"}},
    ]}
    result = _structural_qa(deck_spec)
    assert result["passed"] is True  # None title should not crash or trigger gerundive


def test_format_qa_feedback_structural_only():
    from export import _format_qa_feedback
    result = _format_qa_feedback(
        structural_issues=["Slide count 2 is outside allowed range 3–15"],
        visual_blocking=[],
    )
    assert "QA FEEDBACK" in result
    assert "[STRUCTURAL]" in result
    assert "Slide count 2" in result


def test_format_qa_feedback_visual_only():
    from export import _format_qa_feedback
    result = _format_qa_feedback(
        structural_issues=[],
        visual_blocking=[
            {"slide": 3, "category": "visual", "severity": "blocking",
             "description": "text truncated in body zone"},
        ],
    )
    assert "[VISUAL]" in result
    assert "Slide 3" in result
    assert "text truncated" in result


def test_format_qa_feedback_excludes_warnings():
    from export import _format_qa_feedback
    result = _format_qa_feedback(
        structural_issues=[],
        visual_blocking=[
            {"slide": 2, "category": "quality", "severity": "blocking",
             "description": "gerundive title"},
            {"slide": 4, "category": "coverage", "severity": "warning",
             "description": "agenda item not yet covered"},
        ],
    )
    assert "Slide 2" in result
    assert "Slide 4" not in result  # warnings must not appear


def test_generate_deck_spec_includes_qa_feedback():
    import asyncio, json
    import unittest.mock as mock
    import export as exp

    spec = {"schema_version": 1, "slides": [
        {"layout": "cover", "slots": {"title": "T", "date": "", "duration": ""}},
        {"layout": "bullets", "slots": {"title": "B", "bullets": []}},
        {"layout": "cards-3", "slots": {"title": "C", "cards": [
            {"heading": "A", "content": "a"},
            {"heading": "B", "content": "b"},
            {"heading": "C", "content": "c"},
        ]}},
    ]}
    captured = {}

    async def fake_call(tier, system, user):
        captured["user"] = user
        return json.dumps(spec)

    with mock.patch.object(exp, "_llm_call_slides", fake_call):
        asyncio.run(exp.generate_deck_spec(
            transcript="text",
            recap={},
            instructions=None,
            current_deck_spec=None,
            chain=[{"provider": "test", "model": "m"}],
            qa_feedback="QA FEEDBACK from previous generation — fix these issues:\n[VISUAL] Slide 3: text truncated",
        ))

    assert "QA FEEDBACK" in captured["user"]
    assert "Slide 3" in captured["user"]


def test_format_qa_feedback_empty_returns_empty_string():
    from export import _format_qa_feedback
    result = _format_qa_feedback(structural_issues=[], visual_blocking=[])
    assert result == ""
