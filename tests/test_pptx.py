import io, datetime
import pytest
from pptx import Presentation


TEMPLATE_PATH = "assets/template_cap_pptx.pptx"
THEMES = [
    {
        "name": "Positioning",
        "alignment": ["Clear market leader position", "Enterprise focus agreed"],
        "disagreement": ["SMB vs Enterprise priority split"],
        "unresolved": ["Geographic expansion timing"],
    },
    {
        "name": "Value Proposition",
        "alignment": ["Cost reduction is primary benefit"],
        "disagreement": [],
        "unresolved": ["Secondary benefits not ranked"],
    },
]


def _generate_pptx(themes: list[dict]) -> Presentation:
    """Mirror of the route logic for isolated unit testing."""
    prs = Presentation(TEMPLATE_PATH)

    # Remove all template slides, keeping only the layouts/theme
    _NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(f"{{{_NS_R}}}id")
        if rId:
            prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)

    cover = prs.slides.add_slide(prs.slide_layouts[21])
    cover.placeholders[0].text = "Synthèse Workshop ASE"
    cover.placeholders[22].text = datetime.date.today().strftime("%d/%m/%Y")

    boxes_layout = prs.slide_layouts[35]
    for theme in themes:
        slide = prs.slides.add_slide(boxes_layout)
        slide.placeholders[0].text = theme.get("name", "")

        def fill_box(slide, ph_idx, items, header):
            tf = slide.placeholders[ph_idx].text_frame
            tf.clear()
            tf.text = header
            for item in items:
                p = tf.add_paragraph()
                p.text = f"• {item}"

        fill_box(slide, 22, theme.get("alignment", []), "✓ Alignement")
        fill_box(slide, 35, theme.get("disagreement", []), "✗ Désaccord")
        fill_box(slide, 36, theme.get("unresolved", []), "? Non tranché")

    return prs


def test_pptx_slide_count():
    prs = _generate_pptx(THEMES)
    assert len(prs.slides) == 1 + len(THEMES)


def test_pptx_cover_title():
    prs = _generate_pptx(THEMES)
    cover = prs.slides[0]
    assert cover.placeholders[0].text == "Synthèse Workshop ASE"


def test_pptx_theme_title():
    prs = _generate_pptx(THEMES)
    theme_slide = prs.slides[1]
    assert theme_slide.placeholders[0].text == "Positioning"


def test_pptx_alignment_box():
    prs = _generate_pptx(THEMES)
    theme_slide = prs.slides[1]
    box_text = theme_slide.placeholders[22].text_frame.text
    assert "Alignement" in box_text
    assert "Clear market leader" in box_text


def test_pptx_is_valid_pptx():
    """Check the file can be saved and re-read."""
    prs = _generate_pptx(THEMES)
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    prs2 = Presentation(buf)
    assert len(prs2.slides) == 1 + len(THEMES)


def test_pptx_empty_themes():
    prs = _generate_pptx([])
    assert len(prs.slides) == 1  # cover only
