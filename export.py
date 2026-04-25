#!/usr/bin/env python3
"""
Mimir — Session Graph Export
Generate PDF snapshots and video replays of session knowledge graphs.
Uses Playwright to render the D3 visualization headlessly.

Usage:
    python export.py pdf <session_id>
    python export.py pdf <session_id> -o graph.pdf
    python export.py video <session_id>
    python export.py video <session_id> -o output.mp4 --speed 2.0

Prerequisites:
    pip install playwright && playwright install chromium
"""

import argparse, asyncio, base64, json, os, re, subprocess, sys, tempfile
from datetime import datetime

from log import logger
import aiohttp
import db
import corpus as corpus_module
from prompts.slides import HTML_SLIDES_SYSTEM, deck_spec_system


# ─── LLM config for slides (mirrors routes_facilitator.py pattern) ────────────

ANTHROPIC_API_KEY    = os.environ.get("ANTHROPIC_API_KEY", "")
ANTHROPIC_BASE_URL   = os.environ.get("ANTHROPIC_BASE_URL", "https://api.anthropic.com")
ANTHROPIC_AUTH_TOKEN = os.environ.get("ANTHROPIC_AUTH_TOKEN", "")
HUGIN_BASE_URL       = os.environ.get("HUGIN_BASE_URL", "https://munin.btrbot.com")
HUGIN_CF_ID          = os.environ.get("HUGIN_CF_ID", "")
HUGIN_CF_SECRET      = os.environ.get("HUGIN_CF_SECRET", "")
GEMINI_API_KEY       = os.environ.get("GEMINI_API_KEY", "")
GEMINI_BASE_URL      = "https://generativelanguage.googleapis.com/v1beta/openai"


EXPORT_HTML = os.path.join(os.path.dirname(os.path.abspath(__file__)), "export-graph.html")

_NAV_SCRIPT = """<script>
(function(){
  var slides = document.querySelectorAll('.slide');
  if (!slides.length) return;
  var idx = 0;
  function show(n) {
    slides[idx].classList.remove('active');
    slides[idx].classList.add('exit');
    setTimeout(function(){ slides[idx].classList.remove('exit'); }, 600);
    idx = Math.max(0, Math.min(n, slides.length - 1));
    slides[idx].classList.add('active');
    var dots = document.querySelectorAll('.dot');
    dots.forEach(function(d,i){ d.classList.toggle('active', i===idx); });
    var ctr = document.querySelector('.counter');
    if (ctr) ctr.textContent = (idx+1) + ' / ' + slides.length;
  }
  document.addEventListener('keydown', function(e){
    if (e.key==='ArrowRight'||e.key===' ') show(idx+1);
    if (e.key==='ArrowLeft') show(idx-1);
  });
  document.querySelectorAll('.dot').forEach(function(d,i){
    d.addEventListener('click', function(){ show(i); });
  });
})();
</script>"""


def _inject_nav(html: str) -> str:
    """Inject navigation JS if the LLM didn't generate it (e.g. due to truncation)."""
    if 'ArrowRight' in html or 'keydown' in html:
        return html
    for tag in ('</body>', '</html>'):
        idx = html.lower().rfind(tag)
        if idx >= 0:
            return html[:idx] + _NAV_SCRIPT + '\n' + html[idx:]
    return html + '\n' + _NAV_SCRIPT


def find_peak_snapshot(snapshots: list[dict]) -> dict:
    """Find the snapshot with the most active nodes (best visual)."""
    def score(s):
        nodes = s["graph"].get("nodes", {})
        active = sum(1 for n in nodes.values() if n.get("state") == "active")
        edges = len(s["graph"].get("edges", []))
        return (active, edges)
    return max(snapshots, key=score)


def format_date(ts: float) -> str:
    return datetime.fromtimestamp(ts).strftime("%d %B %Y")


def format_duration(seconds: float) -> str:
    m, s = divmod(int(seconds), 60)
    h, m = divmod(m, 60)
    if h > 0:
        return f"{h}h {m}m"
    return f"{m}m {s}s"


# ─── PDF Export ───

async def export_pdf(session_id: str, output: str, db_path: str = "livemind.db"):
    """Export the peak graph snapshot as a two-page PDF."""
    from playwright.async_api import async_playwright

    if not db._db:
        await db.init_db(db_path)

    # Get session info
    sessions = await db.list_sessions()
    session = next((s for s in sessions if s["id"] == session_id), None)
    if not session:
        sessions = await db.list_sessions(archived=True)
        session = next((s for s in sessions if s["id"] == session_id), None)
    if not session:
        logger.error(f"Session '{session_id}' not found.")
        return

    # Get snapshots
    snapshots = await db.get_session_snapshots(session_id)
    if not snapshots:
        logger.error(f"No snapshots for session '{session_id}'.")
        return

    peak = find_peak_snapshot(snapshots)
    nodes = peak["graph"].get("nodes", {})
    active_count = sum(1 for n in nodes.values() if n.get("state") == "active")
    edge_count = len(peak["graph"].get("edges", []))
    logger.info(f"Peak snapshot: {active_count} nodes, {edge_count} edges")

    # Get summary
    recap_data = await db.get_recap(session_id)
    summary = ""
    if recap_data and recap_data.get("recap"):
        r = recap_data["recap"]
        summary = r.get("elevator_pitch", "") or r.get("summary", "")
    if not summary:
        summary = session.get("summary", "") or ""

    # Compute duration
    duration = ""
    if session.get("ended_at") and session.get("created_at"):
        duration = format_duration(session["ended_at"] - session["created_at"])

    meta = {
        "topic": session.get("topic", "Untitled"),
        "date": format_date(session["created_at"]),
        "summary": summary,
        "stats": {
            "nodes": active_count,
            "edges": edge_count,
            "segments": session.get("segment_count", "—"),
            "duration": duration or "—",
        },
    }

    # Render with Playwright
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page(viewport={"width": 1920, "height": 1080})

        await page.goto("file://" + EXPORT_HTML)

        # Inject data and render
        await page.evaluate(f"renderStatic({json.dumps(peak['graph'])}, {json.dumps(meta)})")
        await page.wait_for_function("window.__READY__ === true", timeout=10000)
        await page.wait_for_timeout(800)  # fonts + render settle

        # Two-page PDF: graph page + metadata page
        await page.pdf(
            path=output,
            width="1920px",
            height="1080px",
            print_background=True,
            margin={"top": "0", "right": "0", "bottom": "0", "left": "0"},
        )

        await browser.close()

    size_kb = os.path.getsize(output) / 1024
    logger.info(f"PDF saved: {output} ({size_kb:.0f} KB)")


# ─── Video Export ───

async def export_video(
    session_id: str,
    output: str,
    db_path: str = "livemind.db",
    speed: float = 1.0,
    max_hold: float = 3.0,
    settle_time: float = 2.5,
):
    """Export session graph evolution as an mp4 video.

    Uses Playwright's built-in screen recording to capture smooth D3 transitions
    instead of jerky frame-by-frame screenshots.
    """
    from playwright.async_api import async_playwright

    if not db._db:
        await db.init_db(db_path)

    snapshots = await db.get_session_snapshots(session_id)
    if not snapshots or len(snapshots) < 2:
        logger.error(f"Not enough snapshots for video (need >= 2, got {len(snapshots or [])}).")
        return

    logger.info(f"Recording {len(snapshots)} snapshots → {output}")
    logger.info(f"Speed: {speed}x, max hold: {max_hold}s, settle: {settle_time}s")

    tmpdir = tempfile.mkdtemp(prefix="mimir-export-")

    try:
        async with async_playwright() as p:
            browser = await p.chromium.launch()

            # Create context with video recording
            context = await browser.new_context(
                viewport={"width": 1920, "height": 1080},
                record_video_dir=tmpdir,
                record_video_size={"width": 1920, "height": 1080},
            )
            page = await context.new_page()

            await page.goto("file://" + EXPORT_HTML)
            await page.evaluate("window.__initVideo__()")
            await page.wait_for_function("window.__READY__ === true", timeout=10000)

            # Brief pause before first snapshot (clean start)
            await page.wait_for_timeout(500)

            for i, snap in enumerate(snapshots):
                # Apply snapshot — D3 animates the transition live
                await page.evaluate(f"window.__applySnapshot__({json.dumps(snap['graph'])})")

                # Let D3 animate and settle naturally
                await page.wait_for_timeout(int(settle_time * 1000))

                # Hold: compressed real-time gap between snapshots
                if i < len(snapshots) - 1:
                    real_gap = snapshots[i + 1]["created_at"] - snap["created_at"]
                    hold_ms = min(max_hold, real_gap / speed) * 1000
                    # Subtract settle time already waited
                    extra_hold = max(0, hold_ms - settle_time * 1000)
                    if extra_hold > 0:
                        await page.wait_for_timeout(int(extra_hold))

                pct = (i + 1) / len(snapshots) * 100
                print(f"\r  Recording: [{pct:5.1f}%] snapshot {i+1}/{len(snapshots)}", end="", flush=True)

            # Brief pause at end
            await page.wait_for_timeout(1500)

            # Close page + context to finalize the recording
            await page.close()
            video_path = await page.video.path()
            await context.close()
            await browser.close()

        logger.info("Converting to mp4...")

        # Convert WebM to MP4 with ffmpeg
        ffmpeg_cmd = [
            "ffmpeg", "-y",
            "-i", video_path,
            "-c:v", "libx264",
            "-pix_fmt", "yuv420p",
            "-preset", "medium",
            "-crf", "20",
            "-loglevel", "error",
            output,
        ]
        result = subprocess.run(ffmpeg_cmd, capture_output=True)
        if result.returncode != 0:
            logger.error(f"ffmpeg error: {result.stderr.decode()}")
            return

        size = os.path.getsize(output)
        size_str = f"{size / 1024 / 1024:.1f} MB" if size > 1024 * 1024 else f"{size / 1024:.0f} KB"
        logger.info(f"Video saved: {output} ({size_str})")

    finally:
        import shutil
        shutil.rmtree(tmpdir, ignore_errors=True)


# ─── HTML Slides Export ───

async def _llm_call_slides(tier: dict, system: str, user: str) -> str:
    """Call a single LLM tier for slides generation. No JSON prefix applied."""
    provider = tier["provider"]
    model = tier["model"]
    timeout = aiohttp.ClientTimeout(total=120)

    if provider == "anthropic":
        url = f"{ANTHROPIC_BASE_URL}/v1/messages"
        headers = {"anthropic-version": "2023-06-01", "content-type": "application/json"}
        if ANTHROPIC_AUTH_TOKEN:
            headers["Authorization"] = f"Bearer {ANTHROPIC_AUTH_TOKEN}"
        else:
            headers["x-api-key"] = ANTHROPIC_API_KEY
        body = {"model": model, "max_tokens": 8192, "system": system,
                "messages": [{"role": "user", "content": user}]}
        async with aiohttp.ClientSession() as s:
            async with s.post(url, headers=headers, json=body, timeout=timeout, ssl=False) as r:
                if r.status != 200:
                    raise RuntimeError(f"Anthropic error {r.status}: {await r.text()}")
                data = await r.json()
                return data["content"][0]["text"]

    if provider == "hugin":
        url = f"{HUGIN_BASE_URL}/v1/chat/completions"
        headers = {"Content-Type": "application/json"}
        if HUGIN_CF_ID and HUGIN_CF_SECRET:
            headers["CF-Access-Client-Id"] = HUGIN_CF_ID
            headers["CF-Access-Client-Secret"] = HUGIN_CF_SECRET
    else:  # gemini
        url = f"{GEMINI_BASE_URL}/chat/completions"
        headers = {"Authorization": f"Bearer {GEMINI_API_KEY}", "Content-Type": "application/json"}

    body = {
        "model": model, "max_tokens": 8192,
        "messages": [{"role": "system", "content": system}, {"role": "user", "content": user}],
    }
    async with aiohttp.ClientSession() as s:
        async with s.post(url, headers=headers, json=body, timeout=timeout, ssl=False) as r:
            if r.status != 200:
                raise RuntimeError(f"{provider} error {r.status}: {await r.text()}")
            data = await r.json()
            return data["choices"][0]["message"]["content"]


async def export_slides(session_id: str, output: str, db_path: str = "livemind.db",
                        chain: list[dict] | None = None):
    """Generate a self-contained HTML slide deck from a session recap via LLM."""
    if not db._db:
        await db.init_db(db_path)

    sessions = await db.list_sessions()
    session = next((s for s in sessions if s["id"] == session_id), None)
    if not session:
        sessions = await db.list_sessions(archived=True)
        session = next((s for s in sessions if s["id"] == session_id), None)
    if not session:
        raise ValueError(f"Session '{session_id}' not found.")

    recap_data = await db.get_recap(session_id)
    if not recap_data or not recap_data.get("recap"):
        raise ValueError(f"Session '{session_id}' has no recap. Generate one first.")
    recap = recap_data["recap"]

    duration = ""
    if session.get("ended_at") and session.get("created_at"):
        duration = format_duration(session["ended_at"] - session["created_at"])

    def _fmt_list(value) -> str:
        if isinstance(value, list):
            return "\n".join(
                f"- {item.get('topics', item) if isinstance(item, dict) else item}"
                + (f": {item['insight']}" if isinstance(item, dict) and item.get("insight") else "")
                for item in value
            )
        return str(value) if value else ""

    user_prompt = f"""Create a complete HTML presentation for this session.

Topic: {session.get('topic', 'Untitled')}
Date: {format_date(session.get('created_at', 0))} | Duration: {duration or 'N/A'}

RECAP:
Pitch: {recap.get('elevator_pitch', '')}
Summary: {recap.get('summary', '')}

Key takeaways:
{_fmt_list(recap.get('key_takeaways', []))}

Non-obvious connections:
{_fmt_list(recap.get('non_obvious_connections', []))}

Tensions & contradictions:
{_fmt_list(recap.get('contradictions', []))}

Output ONLY the complete HTML. Nothing else."""

    if not chain:
        raise RuntimeError("No LLM chain provided. Cannot generate slides.")

    topic = session.get('topic', 'Untitled')
    print(f"Slides export: session={session_id[:8]}… topic='{topic}'")
    print(f"  Prompt: {len(user_prompt)} chars | Chain: {[t['provider']+'/'+t['model'] for t in chain]}")

    html_text = None
    last_error = None
    for tier in chain:
        provider, model = tier["provider"], tier["model"]
        print(f"  Calling {provider}/{model}…", flush=True)
        t0 = __import__("time").time()
        try:
            raw = await _llm_call_slides(tier, HTML_SLIDES_SYSTEM, user_prompt)
            dt = __import__("time").time() - t0
            print(f"  {provider}/{model}: {len(raw)} chars in {dt:.1f}s")
            # Strip markdown fencing if present
            m = re.search(r'```(?:html)?\s*(<!DOCTYPE|<html).*?```', raw, re.DOTALL | re.IGNORECASE)
            if m:
                raw = m.group(0).replace("```html", "").replace("```", "").strip()
            # Find HTML start
            for marker in ("<!DOCTYPE html>", "<!doctype html>", "<html"):
                idx = raw.find(marker)
                if idx >= 0:
                    html_text = raw[idx:]
                    break
            if html_text:
                if not html_text.rstrip().lower().endswith("</html>"):
                    html_text = html_text.rstrip() + "\n</html>"
                print(f"  HTML extracted: {len(html_text)} chars")
                break
            else:
                print(f"  WARNING: no HTML found in response, trying next tier", file=sys.stderr)
        except Exception as e:
            dt = __import__("time").time() - t0
            print(f"  {provider}/{model}: FAILED after {dt:.1f}s — {e}", file=sys.stderr)
            logger.warning("export_slides: tier %s/%s failed: %s", provider, model, e)
            last_error = e

    if not html_text:
        raise RuntimeError(f"All LLM tiers failed. Last error: {last_error}")

    html_text = _inject_nav(html_text)
    os.makedirs(os.path.dirname(os.path.abspath(output)), exist_ok=True)
    with open(output, "w", encoding="utf-8") as f:
        f.write(html_text)

    size_kb = os.path.getsize(output) / 1024
    print(f"Slides saved: {output} ({size_kb:.0f} KB)")


# ─── PPTX Export ───

PPTX_TEMPLATE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets", "template_cap_blank.pptx")


def _fill_bullets(placeholder, items: list[str]) -> None:
    """Fill a text placeholder with bullet items."""
    tf = placeholder.text_frame
    tf.clear()
    if not items:
        return
    tf.text = items[0]
    for item in items[1:]:
        p = tf.add_paragraph()
        p.text = item


def _fmt_recap_items(value) -> list[str]:
    """Flatten recap list items (dicts or strings) to plain strings."""
    if not isinstance(value, list):
        return [str(value)] if value else []
    result = []
    for item in value:
        if isinstance(item, dict):
            topics = item.get("topics", "")
            if isinstance(topics, list):
                topics = " ↔ ".join(topics)
            insight = item.get("insight", "")
            if topics and insight:
                result.append(f"{topics} : {insight}")
            else:
                result.append(topics or insight or str(item))
        else:
            result.append(str(item))
    return [s for s in result if s]


_SKIP_RECAP_KEYS = {"schema_version", "transcript_stats"}
_V3_STRUCTURED_KEYS = {
    "positioning", "value_proposition",
    "positioning_statement", "scope_boundaries_non_goals",
}
_POSITIONING_LABELS = {
    "what_to_sell": "What to sell?",
    "why_now": "Why now?",
    "why_well_positioned": "Why are we well positioned?",
    "to_whom": "To whom?",
}
_VP_LABELS = {
    "what_we_do": "What do we do?",
    "how_we_do_it": "How do we do it?",
    "how_we_get_paid": "How do we get paid?",
}


def _format_recap(recap: dict) -> str:
    """Convert a recap dict to human-readable structured text for LLM prompts."""
    lines = []

    positioning = recap.get("positioning")
    if positioning:
        lines.append("=== POSITIONING ===")
        for key, label in _POSITIONING_LABELS.items():
            val = positioning.get(key)
            if val:
                items_str = "; ".join(str(v) for v in val) if isinstance(val, list) else str(val)
                lines.append(f"  {label:<38} → {items_str}")
        lines.append("")

    value_prop = recap.get("value_proposition")
    if value_prop:
        lines.append("=== VALUE PROPOSITION ===")
        for key, label in _VP_LABELS.items():
            val = value_prop.get(key)
            if val:
                items_str = "; ".join(str(v) for v in val) if isinstance(val, list) else str(val)
                lines.append(f"  {label:<38} → {items_str}")
        lines.append("")

    pos_stmt = recap.get("positioning_statement")
    if pos_stmt:
        lines.append("=== POSITIONING STATEMENT ===")
        lines.append(f'"{pos_stmt}"')
        lines.append("")

    scope = recap.get("scope_boundaries_non_goals")
    if scope:
        lines.append("=== SCOPE / BOUNDARIES / NON-GOALS ===")
        items = scope if isinstance(scope, list) else [scope]
        for item in items:
            lines.append(f"  - {item}")
        lines.append("")

    # Unknown / V2-style keys
    _known = _SKIP_RECAP_KEYS | _V3_STRUCTURED_KEYS
    for key, val in recap.items():
        if key in _known or not val:
            continue
        label = key.replace("_", " ").upper()
        lines.append(f"=== {label} ===")
        if isinstance(val, list):
            for item in val:
                if isinstance(item, dict):
                    topics = item.get("topics", "")
                    if isinstance(topics, list):
                        topics = " ↔ ".join(topics)
                    insight = item.get("insight", "")
                    line = f"{topics} : {insight}" if topics and insight else (topics or insight or str(item))
                    lines.append(f"  - {line}")
                else:
                    lines.append(f"  - {item}")
        elif isinstance(val, dict):
            lines.append(f"  {json.dumps(val, ensure_ascii=False)}")
        else:
            lines.append(f"  {val}")
        lines.append("")

    return "\n".join(lines).strip()


def _format_corpus_for_slides(docs: list[dict], max_chars: int = 300_000) -> str:
    """Format active corpus docs for injection into the slides generation prompt.
    Groups chunks by source document with metadata headers.
    If total content exceeds max_chars, includes only metadata headers (no chunk content).
    """
    if not docs:
        return ""

    from collections import defaultdict
    by_source: dict[str, list[str]] = defaultdict(list)
    meta_by_source: dict[str, dict] = {}
    for doc in docs:
        src = doc["source"] or doc["title"]
        by_source[src].append(doc["content"])
        if src not in meta_by_source:
            meta_by_source[src] = {
                "label": doc.get("label") or src,
                "role": doc.get("role"),
                "key_messages": doc.get("key_messages"),
                "usages": doc.get("usages"),
            }

    def _build_header(m: dict) -> str:
        header = f"=== {m['label']} ==="
        if m["role"]:
            header += f"\nRôle : {m['role']}"
        if m["key_messages"]:
            try:
                kms = json.loads(m["key_messages"]) if isinstance(m["key_messages"], str) else (m["key_messages"] or [])
            except Exception:
                kms = []
            if kms:
                header += "\nMessages clés :\n" + "\n".join(f"  • {k}" for k in kms)
        if m["usages"]:
            try:
                us = json.loads(m["usages"]) if isinstance(m["usages"], str) else (m["usages"] or [])
            except Exception:
                us = []
            if us:
                header += "\nUsages : " + " · ".join(us)
        return header

    # Build full sections (header + content)
    sections = []
    for src, chunks in by_source.items():
        header = _build_header(meta_by_source[src]) + "\n---"
        sections.append(header + "\n" + "\n\n".join(chunks))

    full = "\n\n".join(sections)
    if len(full) <= max_chars:
        return full

    # Fallback: headers only (no chunk content) to stay within budget
    logger.warning(f"Corpus too large for slides ({len(full)} chars > {max_chars}), including metadata only")
    headers_only = []
    for src in by_source:
        headers_only.append(_build_header(meta_by_source[src]))
    return "\n\n".join(headers_only)


def _build_user_prompt(
    recap: dict,
    transcript: str,
    instructions: str | None,
    current_deck_spec: dict | None,
    session_topic: str = "",
    session_date: str = "",
    reminder: str = "",
    corpus_block: str = "",
    qa_feedback: str = "",
) -> str:
    """Build the LLM user prompt for deck_spec generation."""
    parts = []
    if qa_feedback and qa_feedback.strip():
        parts.append(qa_feedback.strip())
    if instructions and instructions.strip():
        parts.append(
            "INSTRUCTIONS (apply to ALL slides — "
            "titles, content, language, tone):\n" + instructions.strip()
        )
    if session_topic or session_date:
        ctx = []
        if session_topic:
            ctx.append(f"Topic: {session_topic}")
        if session_date:
            ctx.append(f"Date: {session_date}")
        parts.append("CONTEXT:\n" + "\n".join(ctx))
    parts.append("RECAP:\n" + _format_recap(recap))
    parts.append("TRANSCRIPT (excerpt):\n" + transcript)
    if corpus_block:
        parts.append(
            "REFERENCE DOCUMENTS AVAILABLE\n\n"
            "These documents are supplementary context only.\n"
            "ABSOLUTE RULE: the session transcript is the primary source for slide content.\n"
            "Do not introduce any information not discussed during the session.\n\n"
            "However:\n"
            "- If a topic from the transcript resonates with a corpus document, you may enrich\n"
            "  the slide content with specific elements from that document.\n"
            "- If a document contains relevant elements NOT covered in the transcript,\n"
            "  you may add to the slide's \"notes\" field:\n"
            "  \"\\u2139\\ufe0f [Document name] contains additional elements on this topic.\"\n"
            "- Always cite the source document in brackets when borrowing an element.\n\n"
            + corpus_block
        )
    if current_deck_spec:
        parts.append(
            "CURRENT DECK:\n"
            + json.dumps(current_deck_spec, ensure_ascii=False, indent=2)
        )
    if reminder:
        parts.append(reminder)
    parts.append("Generate the complete deck_spec JSON.")
    return "\n\n".join(parts)


_CARDS_HEADING_MARKER = "Lorem ipsum dolor"
_CARDS_CONTENT_MARKER = "Aenean vulputate"


def _fill_cards_shapes(slide, cards: list[dict]) -> None:
    """Replace heading/content placeholder text in a copied cards slide.

    Identifies shapes by their placeholder text (Lorem ipsum → heading,
    Aenean vulputate → content), sorts by vertical then horizontal position,
    and fills them in order from the cards list.
    """
    headings = []
    contents = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text
        if _CARDS_HEADING_MARKER in txt:
            headings.append(shape)
        elif _CARDS_CONTENT_MARKER in txt:
            contents.append(shape)

    headings.sort(key=lambda s: (s.top, s.left))
    contents.sort(key=lambda s: (s.top, s.left))

    for i, card in enumerate(cards):
        if i < len(headings):
            headings[i].text_frame.text = card.get("heading", "")
        if i < len(contents):
            contents[i].text_frame.text = card.get("content", "")


def _copy_slide_from_template(
    prs,
    tmpl,
    slide_idx: int,
    title: str,
    cards: list[dict],
) -> None:
    """Add a slide to prs by copying slide[slide_idx] from tmpl, then filling content.

    Copies the full spTree XML and all OPC relationships (images, OLE objects, tags)
    from the source slide, remapping rIds to avoid dangling references in the new slide.

    Args:
        prs: target Presentation (being assembled)
        tmpl: source Presentation (template loaded separately)
        slide_idx: 0-based index of the template slide to copy
        title: slide title text
        cards: list of {"heading": str, "content": str} dicts
    """
    import copy as _copy
    import re
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    src = tmpl.slides[slide_idx]
    src_part = src.part

    # Find matching layout in prs by name
    src_layout_name = src.slide_layout.name
    matching_layout = next(
        (l for l in prs.slide_layouts if l.name == src_layout_name),
        prs.slide_layouts[35],
    )
    new_slide = prs.slides.add_slide(matching_layout)
    new_part = new_slide.part

    # Copy all non-layout, non-tag relationships from src to new_slide, building rId remap.
    # Reuse parts already present in prs (same partname) to avoid duplicate ZIP entries.
    # Tags (think-cell slide metadata) are slide-specific and must not be shared between
    # slides; skipping them here prevents PPTX corruption.
    _skip_reltypes = {
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/tags",
    }
    # Build a map of existing parts in prs to avoid duplicate ZIP entries
    prs_parts = {str(p.partname): p for p in new_part.package.iter_parts()}
    rId_map = {}
    for rId, rel in src_part.rels.items():
        if rel.reltype in _skip_reltypes:
            continue  # layout already set by add_slide; tags are slide-specific
        try:
            # Prefer an existing part at the same path to avoid duplicate ZIP entries
            target = prs_parts.get(str(rel.target_part.partname), rel.target_part)
            new_rId = new_part.relate_to(target, rel.reltype)
            rId_map[rId] = new_rId
        except Exception:
            # Skip relationships that can't be copied (e.g. external refs)
            pass

    # Deep-copy the spTree XML from source
    sp_tree_copy = _copy.deepcopy(src.shapes._spTree)

    # Remove <p:tags> elements from <p:custDataLst> within the spTree copy.
    # These are think-cell identifiers stored inside shape custom-data; they reference
    # the tag relationship (now skipped above) and would create dangling rId refs.
    _P_CUSTDATALST = '{http://schemas.openxmlformats.org/presentationml/2006/main}custDataLst'
    _P_TAGS = '{http://schemas.openxmlformats.org/presentationml/2006/main}tags'
    for cust in list(sp_tree_copy.iter(_P_CUSTDATALST)):
        for tag_elem in list(cust.findall(_P_TAGS)):
            cust.remove(tag_elem)

    # Remap all r:id and r:embed attribute values using rId_map
    if rId_map:
        xml_bytes = sp_tree_copy.xml.encode("utf-8")
        for old_rId, new_rId in rId_map.items():
            xml_bytes = xml_bytes.replace(
                f'"{old_rId}"'.encode(),
                f'"{new_rId}"'.encode(),
            )
        from lxml import etree
        sp_tree_copy = etree.fromstring(xml_bytes)

    # Replace the new slide's spTree with the remapped copy
    sp_tree = new_slide.shapes._spTree
    for child in list(sp_tree):
        sp_tree.remove(child)
    for child in sp_tree_copy:
        sp_tree.append(_copy.deepcopy(child))

    # Update title
    try:
        new_slide.placeholders[0].text = title
    except (KeyError, IndexError):
        logger.warning(f"_copy_slide_from_template: placeholder[0] not found in slide {slide_idx}")

    # Fill cards content
    _fill_cards_shapes(new_slide, cards)


# ─── Layout catalogue (exposed to LLM) ───────────────────────────────────────

LAYOUT_CATALOG = {
    "cover": {
        "description": "Title slide",
        "slots": "title (str), date (str), duration (str)",
        "layout_idx": 0,
    },
    "text-large": {
        "description": "Title + large body text",
        "slots": "title (str), body (str)",
        "layout_idx": 21,
    },
    "quote-large": {
        "description": "Featured quote or pitch statement",
        "slots": "title (str), body (str)",
        "layout_idx": 21,
    },
    "bullets": {
        "description": "Title + bullet list (max 6 items, max 15 words each)",
        "slots": "title (str), bullets (list[str])",
        "layout_idx": 21,
    },
    "three-columns": {
        "description": "Three equal columns",
        "slots": "title (str), col1 (str), col2 (str), col3 (str)",
        "layout_idx": 35,
    },
    "two-columns": {
        "description": "Two columns",
        "slots": "title (str), left (str), right (str)",
        "layout_idx": 35,
    },
    "concepts": {
        "description": "Term cloud + relationships",
        "slots": "title (str), terms (list[str]), edges (list[str])",
        "layout_idx": 48,
    },
    "cards-3": {
        "description": "3 structured cards side by side (heading + content)",
        "slots": 'title (str), cards (list of {"heading": str, "content": str}, exactly 3 items)',
        "layout_idx": 35,
        "slide_copy_idx": 8,   # slide 9 du template, 0-indexé
    },
    "cards-4": {
        "description": "4 structured cards (heading + content)",
        "slots": 'title (str), cards (list of {"heading": str, "content": str}, exactly 4 items)',
        "layout_idx": 35,
        "slide_copy_idx": 9,   # slide 10 du template
    },
    "cards-5": {
        "description": "5 structured cards (heading + content)",
        "slots": 'title (str), cards (list of {"heading": str, "content": str}, exactly 5 items)',
        "layout_idx": 35,
        "slide_copy_idx": 10,  # slide 11 du template
    },
    "cards-4-rounded": {
        "description": "4 rounded-style cards (heading + content, visual variant of cards-4)",
        "slots": 'title (str), cards (list of {"heading": str, "content": str}, exactly 4 items)',
        "layout_idx": 35,
        "slide_copy_idx": 11,  # slide 12 du template
    },
    "divider": {
        "description": "Section separator (number + section title)",
        "slots": "title (str), number (str)",
        "layout_idx": 14,  # Divider 2 in template_cap_blank.pptx
    },
}

_LAYOUT_CATALOG_STR = "\n".join(
    f'- "{name}": {info["description"]}. Slots: {info["slots"]}'
    for name, info in LAYOUT_CATALOG.items()
)


def _structural_qa(deck_spec: dict) -> dict:
    """Check deck_spec structure before assembly. Returns {"passed": bool, "issues": list[str]}."""
    issues = []
    slides = deck_spec.get("slides", [])

    if not (3 <= len(slides) <= 15):
        issues.append(f"Slide count {len(slides)} is outside allowed range 3–15")

    consecutive = 1
    for i in range(1, len(slides)):
        if slides[i].get("layout") == slides[i - 1].get("layout"):
            consecutive += 1
            if consecutive > 3:
                issues.append(
                    f"Slides {i - 2}–{i + 1}: {consecutive} consecutive "
                    f"'{slides[i]['layout']}' layouts (max 3)"
                )
                break
        else:
            consecutive = 1

    gerundive = re.compile(r"^[A-Z][a-zA-Z]+ing\b")
    bad_titles = [
        f"Slide {i + 1}: \"{s.get('slots', {}).get('title', '')}\""
        for i, s in enumerate(slides)
        if gerundive.match(s.get("slots", {}).get("title") or "")
    ]
    if len(bad_titles) > 2:
        issues.append(
            f"{len(bad_titles)} gerundive slide titles (max 2): {'; '.join(bad_titles)}"
        )

    return {"passed": len(issues) == 0, "issues": issues}


def _format_qa_feedback(structural_issues: list, visual_blocking: list) -> str:
    """Format blocking QA issues as a feedback block for the next generate_deck_spec call.

    Returns empty string when there are no blocking issues.
    """
    blocking_visual = [i for i in visual_blocking if i.get("severity") != "warning"]
    if not structural_issues and not blocking_visual:
        return ""
    lines = ["QA FEEDBACK from previous generation — fix these issues:"]
    for issue in structural_issues:
        lines.append(f"[STRUCTURAL] {issue}")
    for issue in blocking_visual:
        cat = issue.get("category", "issue").upper()
        slide = issue.get("slide", "?")
        desc = issue.get("description", "")
        lines.append(f"[{cat}] Slide {slide}: {desc}")
    return "\n".join(lines)


def _soffice_path() -> str | None:
    """Find LibreOffice soffice binary on macOS or Linux. Returns None if not found."""
    import shutil
    mac_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    if os.path.exists(mac_path):
        return mac_path
    return shutil.which("soffice")


def _png_slide_index(path: str) -> int:
    """Return the trailing integer from a PNG filename for numeric sort ordering."""
    m = re.search(r"(\d+)\.png$", path)
    return int(m.group(1)) if m else 0


def _pptx_to_thumbnails(pptx_path: str, tmpdir: str) -> list[str]:
    """Convert a .pptx to per-slide PNG thumbnails via LibreOffice headless.

    Returns sorted list of PNG paths. Returns [] if soffice is unavailable or fails.
    """
    soffice = _soffice_path()
    if not soffice:
        return []
    try:
        result = subprocess.run(
            [soffice, "--headless", "--convert-to", "png", "--outdir", tmpdir, pptx_path],
            capture_output=True, timeout=60,
        )
        if result.returncode != 0:
            logger.error(f"_pptx_to_thumbnails: soffice failed: {result.stderr.decode()}")
            return []
    except Exception as e:
        logger.error(f"_pptx_to_thumbnails: soffice error: {e}")
        return []
    pngs = sorted(
        [os.path.join(tmpdir, f) for f in os.listdir(tmpdir) if f.lower().endswith(".png")],
        key=_png_slide_index,
    )
    return pngs


_VISUAL_QA_PROMPT = """\
You are inspecting an auto-generated PowerPoint presentation.
Inspect each slide and evaluate the following criteria:

1. CAPGEMINI TEMPLATE COMPLIANCE
   Each slide must use the Capgemini Invent dark navy background,
   corporate blue (#0058AB), and must not display a white or generic background.

2. LAYOUT DIVERSITY
   The deck must alternate between layouts (cards, bullets, columns, quotes).
   Flag if more than 3 consecutive slides share the same visual structure.

3. EXECUTIVE COMMITTEE CONTENT QUALITY
   For each content slide (excluding cover and dividers):
   - Is the title a DECLARATIVE ASSERTION (not a gerundive description)?
     Good example: "The Transform-then-Run model is our primary competitive moat"
     Bad example: "Positioning Capgemini Invent in the Market"
   - Does the content name a decision, implication, risk, or action?
   - Is the language appropriate for a Group Executive Committee audience?

4. AGENDA COVERAGE (informational only — never blocking)
   Note which of the following themes are present in the deck.
   Absence is expected if the topic has not yet been discussed in the live session.
   Positioning: What we sell (Transform vs run; Asset/IP-led; Front/Core/Back) /
     Why now (Market inflexion; IOPs momentum; Agentic operations) /
     Why well positioned (Tri-pod; Orchestrator; Credibility) /
     To Whom (CXO play; Customer tier; Archetypes) / Position statement
   Value Prop: What we do (Value engine; E2E reinvention) /
     How we do it (Deals anatomy; Capabilities aggregation) /
     What we get paid for (Value/Risks/Cash; Shared accountability)
   Targets & horizon / Priorities & orchestration / Scope & non-goals

5. VISUAL STANDARDS
   - Text truncated or overflowing its shape or zone
   - Overlapping elements
   - Visible unfilled placeholder ("Lorem ipsum", unexpected empty zone)

Reply ONLY in strict JSON (no markdown fences, no extra text):
{"passed": true, "issues": [{"slide": 1, "category": "template", "severity": "blocking", "description": "..."}], "summary": "..."}
"passed" must be false if at least one issue has severity "blocking".\
"""


def _parse_visual_qa_response(raw: str) -> dict:
    """Parse the JSON string returned by the visual QA Claude call.

    Returns {"passed": bool, "issues": list[dict], "warnings": list[dict], "summary": str}.
    Blocking issues go into "issues"; warning issues go into "warnings".
    """
    raw = raw.strip()
    if raw.startswith("```"):
        raw = re.sub(r"^```[a-zA-Z]*\n?", "", raw)
        raw = re.sub(r"\n?```$", "", raw)
        raw = raw.strip()
    data = json.loads(raw)
    issues = data.get("issues", [])
    blocking = [i for i in issues if i.get("severity") == "blocking"]
    warnings = [i for i in issues if i.get("severity") == "warning"]
    return {
        "passed": len(blocking) == 0,
        "issues": blocking,
        "warnings": warnings,
        "summary": data.get("summary", ""),
    }


async def _visual_qa(
    pptx_path: str,
    deck_spec: dict,
    session_topic: str = "",
) -> dict:
    """Run visual QA on a .pptx via LibreOffice thumbnails + Claude Haiku vision.

    Returns {"passed": bool, "issues": list[dict], "warnings": list[dict], "summary": str}.
    Gracefully skips (passed=True) if LibreOffice or Anthropic credentials are absent.
    """
    _skip = {"passed": True, "issues": [], "warnings": [], "summary": "Visual QA skipped"}

    if not ANTHROPIC_API_KEY and not ANTHROPIC_AUTH_TOKEN:
        logger.warning("_visual_qa: no Anthropic credentials — skipping")
        return {**_skip, "summary": "Visual QA skipped (no Anthropic credentials)"}

    if not _soffice_path():
        logger.warning("_visual_qa: LibreOffice not found — skipping")
        return {**_skip, "summary": "Visual QA skipped (LibreOffice not found)"}

    with tempfile.TemporaryDirectory(prefix="mimir-qa-") as tmpdir:
        png_paths = _pptx_to_thumbnails(pptx_path, tmpdir)
        if not png_paths:
            logger.warning("_visual_qa: thumbnail generation failed — skipping")
            return {**_skip, "summary": "Visual QA skipped (thumbnail generation failed)"}

        content = []

        # Optional: inject one reference slide for style guidance
        ref_deck = os.environ.get("PPTX_REFERENCE_DECK", "")
        if ref_deck and os.path.exists(ref_deck):
            with tempfile.TemporaryDirectory(prefix="mimir-ref-") as ref_tmpdir:
                ref_pngs = _pptx_to_thumbnails(ref_deck, ref_tmpdir)
                if ref_pngs:
                    with open(ref_pngs[0], "rb") as f:
                        ref_b64 = base64.b64encode(f.read()).decode()
                    content.append({"type": "text", "text": "Reference: this is a correctly formatted Capgemini slide:"})
                    content.append({"type": "image", "source": {"type": "base64", "media_type": "image/png", "data": ref_b64}})

        for i, path in enumerate(png_paths):
            with open(path, "rb") as f:
                img_b64 = base64.b64encode(f.read()).decode()
            content.append({"type": "text", "text": f"Slide {i + 1}:"})
            content.append({"type": "image", "source": {"type": "base64", "media_type": "image/png", "data": img_b64}})

        ctx_note = "on the topic of Intelligent Operations (IOPS)"
        if session_topic:
            ctx_note += f" — {session_topic}"
        context_header = (
            f"Context: This deck is intended for the Capgemini Group Executive Committee, {ctx_note}.\n\n"
        )
        content.append({"type": "text", "text": context_header + _VISUAL_QA_PROMPT})

        url = f"{ANTHROPIC_BASE_URL}/v1/messages"
        headers = {"anthropic-version": "2023-06-01", "content-type": "application/json"}
        if ANTHROPIC_AUTH_TOKEN:
            headers["Authorization"] = f"Bearer {ANTHROPIC_AUTH_TOKEN}"
        else:
            headers["x-api-key"] = ANTHROPIC_API_KEY

        body = {
            "model": "claude-haiku-4-5-20251001",
            "max_tokens": 1024,
            "messages": [{"role": "user", "content": content}],
        }

        timeout = aiohttp.ClientTimeout(total=120)
        try:
            async with aiohttp.ClientSession() as session:
                async with session.post(url, headers=headers, json=body, timeout=timeout, ssl=False) as r:
                    if r.status != 200:
                        text = await r.text()
                        logger.error(f"_visual_qa: API error {r.status}: {text[:200]}")
                        return {**_skip, "summary": f"Visual QA skipped (API error {r.status})"}
                    data = await r.json()
                    raw = data["content"][0]["text"]
        except Exception as e:
            logger.error(f"_visual_qa: request failed: {e}")
            return {**_skip, "summary": f"Visual QA skipped ({e})"}

    try:
        return _parse_visual_qa_response(raw)
    except (json.JSONDecodeError, KeyError) as e:
        logger.error(f"_visual_qa: JSON parse error: {e} — raw: {raw[:200]}")
        return {**_skip, "summary": "Visual QA skipped (invalid JSON response)"}


def _clear_slides(prs) -> None:
    """Remove all existing slides from a Presentation, keeping slide layouts and masters."""
    from pptx.oxml.ns import qn
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn("r:id"))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def _assemble_pptx(deck_spec: dict, output: str) -> None:
    """Assemble a .pptx file deterministically from a deck_spec dict."""
    from pptx import Presentation

    prs = Presentation(PPTX_TEMPLATE)
    _clear_slides(prs)

    _COPY_SLIDE_LAYOUTS = {"cards-3", "cards-4", "cards-5", "cards-4-rounded"}
    tmpl = Presentation(PPTX_TEMPLATE)  # source for copy-slide layouts

    for slide_def in deck_spec.get("slides", []):
        layout_name = slide_def.get("layout", "bullets")
        slots = slide_def.get("slots", {})

        # Fallback for unknown layouts
        if layout_name not in LAYOUT_CATALOG:
            logger.warning(f"_assemble_pptx: unknown layout '{layout_name}', falling back to 'bullets'")
            layout_name = "bullets"

        # Copy-slide layouts: copy entire slide from template, fill content
        if layout_name in _COPY_SLIDE_LAYOUTS:
            slide_copy_idx = LAYOUT_CATALOG[layout_name]["slide_copy_idx"]
            _copy_slide_from_template(
                prs,
                tmpl,
                slide_idx=slide_copy_idx,
                title=slots.get("title", ""),
                cards=slots.get("cards", []),
            )
            continue

        # Layout-based slides: add new slide from layout
        layout_idx = LAYOUT_CATALOG[layout_name]["layout_idx"]
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        if layout_name == "cover":
            slide.placeholders[0].text = slots.get("title", "")
            slide.placeholders[10].text = slots.get("date", "")
            if slots.get("duration"):
                slide.placeholders[11].text = slots["duration"]

        elif layout_name == "divider":
            slide.placeholders[0].text = slots.get("title", "")
            try:
                slide.placeholders[23].text = slots.get("number", "")  # ph[23] = body text in Divider 2
            except (KeyError, IndexError):
                logger.warning("_assemble_pptx: divider ph[23] not found")

        elif layout_name in ("text-large", "quote-large"):
            slide.placeholders[0].text = slots.get("title", "")
            slide.placeholders[22].text = slots.get("body", "")

        elif layout_name == "bullets":
            slide.placeholders[0].text = slots.get("title", "")
            items = slots.get("bullets", [])
            if isinstance(items, str):
                items = [items]
            _fill_bullets(slide.placeholders[22], items)

        elif layout_name == "three-columns":
            slide.placeholders[0].text = slots.get("title", "")
            slide.placeholders[22].text = slots.get("col1", "")
            slide.placeholders[35].text = slots.get("col2", "")
            slide.placeholders[36].text = slots.get("col3", "")

        elif layout_name == "two-columns":
            slide.placeholders[0].text = slots.get("title", "")
            slide.placeholders[22].text = slots.get("left", "")
            slide.placeholders[35].text = slots.get("right", "")

        elif layout_name == "concepts":
            slide.placeholders[0].text = slots.get("title", "")
            tf = slide.placeholders[13].text_frame
            tf.clear()
            terms = slots.get("terms", [])
            edges = slots.get("edges", [])
            tf.text = ", ".join(terms[:24])
            if edges:
                tf.add_paragraph().text = ""
                p = tf.add_paragraph()
                p.text = "Relations :"
                for edge in edges[:10]:
                    tf.add_paragraph().text = f"  {edge}"

    os.makedirs(os.path.dirname(os.path.abspath(output)), exist_ok=True)
    prs.save(output)
    size_kb = os.path.getsize(output) / 1024
    logger.info(f"PPTX saved: {output} ({size_kb:.0f} KB)")


async def generate_deck_spec(
    transcript: str,
    recap: dict,
    instructions: str | None,
    current_deck_spec: dict | None,
    chain: list[dict],
    session_topic: str = "",
    session_date: str = "",
    corpus_block: str = "",
    qa_feedback: str = "",
) -> dict:
    """Generate a deck_spec from transcript + recap + instructions via LLM.

    Returns the deck_spec dict. Retries once on invalid JSON.
    Raises RuntimeError if all tiers fail.
    """
    # Truncate transcript — keep tail (most recent content is most relevant)
    max_transcript = 200000
    if len(transcript) > max_transcript:
        transcript = "[...]\n" + transcript[-max_transcript:]

    last_error = None
    for attempt in range(2):
        reminder = (
            "" if attempt == 0
            else "REMINDER: output JSON only, no text, no markdown fences."
        )
        user_prompt = _build_user_prompt(
            recap=recap,
            transcript=transcript,
            instructions=instructions,
            current_deck_spec=current_deck_spec,
            session_topic=session_topic,
            session_date=session_date,
            reminder=reminder,
            corpus_block=corpus_block,
            qa_feedback=qa_feedback,
        )

        for tier in chain:
            provider, model = tier["provider"], tier["model"]
            try:
                raw = await _llm_call_slides(tier, deck_spec_system(_LAYOUT_CATALOG_STR), user_prompt)
                raw = raw.strip()
                if raw.startswith("```"):
                    raw = re.sub(r"^```[a-zA-Z]*\n?", "", raw)
                    raw = re.sub(r"\n?```$", "", raw)
                    raw = raw.strip()
                spec = json.loads(raw)
                if len(spec.get("slides", [])) > 15:
                    logger.warning(
                        f"generate_deck_spec: {len(spec['slides'])} slides, truncating to 15"
                    )
                    spec["slides"] = spec["slides"][:15]
                return spec
            except json.JSONDecodeError as e:
                logger.warning(f"generate_deck_spec attempt {attempt+1}: JSON parse error — {e}")
                last_error = e
                break  # retry outer loop with reminder
            except Exception as e:
                logger.warning(f"generate_deck_spec tier {provider}/{model} failed: {e}")
                last_error = e
                continue  # try next tier

    raise RuntimeError(f"generate_deck_spec: all attempts failed. Last error: {last_error}")


async def export_pptx(session_id: str, output: str, db_path: str = "livemind.db",
                      chain: list[dict] | None = None) -> None:
    """Generate a PPTX slide deck via LLM (deck_spec) then assemble deterministically."""
    if not chain:
        raise RuntimeError("export_pptx requires a chain. Pass chain= from app.py.")

    if not db._db:
        await db.init_db(db_path)

    # Load session
    sessions = await db.list_sessions()
    session = next((s for s in sessions if s["id"] == session_id), None)
    if not session:
        sessions = await db.list_sessions(archived=True)
        session = next((s for s in sessions if s["id"] == session_id), None)
    if not session:
        raise ValueError(f"Session '{session_id}' not found.")

    # Load recap (required)
    recap_data = await db.get_recap(session_id)
    if not recap_data or not recap_data.get("recap"):
        raise ValueError(f"Session '{session_id}' has no recap. Generate one first.")
    recap = recap_data["recap"]

    # Load pptx_data (instructions + existing deck_spec)
    pptx_data = await db.get_pptx_data(session_id)
    instructions = pptx_data["instructions"] if pptx_data else None
    current_deck_spec = pptx_data["deck_spec"] if pptx_data else None

    # Load transcript
    segments = await db.get_session_transcript(session_id)
    transcript = " ".join(s["text"] for s in segments if not s.get("is_partial"))

    # Load active corpus docs for contextual enrichment
    corpus_docs = await corpus_module.get_active_docs(db._db) if db._db else []
    corpus_block = _format_corpus_for_slides(corpus_docs)

    topic = session.get("topic", "Untitled")
    logger.info(f"PPTX export: session={session_id[:8]}… topic='{topic}' "
          f"transcript={len(transcript)} chars instructions={'yes' if instructions else 'no'} "
          f"deck_spec={'update' if current_deck_spec else 'new'} "
          f"corpus={len(corpus_docs)} chunks")

    try:
        session_date = datetime.fromtimestamp(float(session.get("created_at", 0))).strftime("%d %B %Y")
    except Exception:
        session_date = ""

    # Generate deck_spec via LLM
    deck_spec = await generate_deck_spec(
        transcript=transcript,
        recap=recap,
        instructions=instructions,
        current_deck_spec=current_deck_spec,
        chain=chain,
        session_topic=topic,
        session_date=session_date,
        corpus_block=corpus_block,
    )

    # Persist deck_spec
    served_model = chain[0]["model"]
    await db.save_deck_spec(session_id, deck_spec, served_model)

    # Assemble PPTX
    _assemble_pptx(deck_spec, output)


# ─── CLI ───

def main():
    parser = argparse.ArgumentParser(
        description="Export Mimir session graphs as PDF or video"
    )
    sub = parser.add_subparsers(dest="command", required=True)

    pdf_parser = sub.add_parser("pdf", help="Export peak graph snapshot as PDF")
    pdf_parser.add_argument("session_id", help="Session ID to export")
    pdf_parser.add_argument("-o", "--output", help="Output file path (default: <session_id>.pdf)")
    pdf_parser.add_argument("--db", default="livemind.db", help="Database path")

    vid_parser = sub.add_parser("video", help="Export graph evolution as mp4 video")
    vid_parser.add_argument("session_id", help="Session ID to export")
    vid_parser.add_argument("-o", "--output", help="Output file path (default: <session_id>.mp4)")
    vid_parser.add_argument("--speed", type=float, default=2.0, help="Playback speed multiplier (default: 2.0)")
    vid_parser.add_argument("--max-hold", type=float, default=3.0, help="Max seconds per snapshot (default: 3.0)")
    vid_parser.add_argument("--settle", type=float, default=2.5, help="Seconds for D3 animation to settle per snapshot (default: 2.5)")
    vid_parser.add_argument("--db", default="livemind.db", help="Database path")

    args = parser.parse_args()

    if args.command == "pdf":
        output = args.output or f"{args.session_id}.pdf"
        asyncio.run(export_pdf(args.session_id, output, args.db))

    elif args.command == "video":
        output = args.output or f"{args.session_id}.mp4"
        asyncio.run(export_video(
            args.session_id, output, args.db,
            speed=args.speed, max_hold=args.max_hold, settle_time=args.settle,
        ))


if __name__ == "__main__":
    main()
