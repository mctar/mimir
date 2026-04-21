"""
FastAPI routes for ASE Workshop Facilitator view.
Mount via: app.include_router(routes_facilitator.router)
Configure via: routes_facilitator.configure(...)
"""

import json, re, time, io, os, datetime
import aiohttp
from fastapi import APIRouter, Request
from fastapi.responses import FileResponse, JSONResponse, StreamingResponse

import db as db_module
import corpus as corpus_module

router = APIRouter()

# ─── State injected by app.py at startup ────────────────────────────────────
_get_session_id = lambda: None
_broadcast = None
_get_llm_chain = lambda: []
_get_db_conn = lambda: None  # returns the shared aiosqlite.Connection (db._db)

# ─── Provider config (mirrors app.py) ────────────────────────────────────────
ANTHROPIC_API_KEY    = os.environ.get("ANTHROPIC_API_KEY", "")
ANTHROPIC_BASE_URL   = os.environ.get("ANTHROPIC_BASE_URL", "https://api.anthropic.com")
ANTHROPIC_AUTH_TOKEN = os.environ.get("ANTHROPIC_AUTH_TOKEN", "")
HUGIN_BASE_URL       = os.environ.get("HUGIN_BASE_URL", "https://munin.btrbot.com")
HUGIN_CF_ID          = os.environ.get("HUGIN_CF_ID", "")
HUGIN_CF_SECRET      = os.environ.get("HUGIN_CF_SECRET", "")
GEMINI_API_KEY       = os.environ.get("GEMINI_API_KEY", "")
GEMINI_BASE_URL      = "https://generativelanguage.googleapis.com/v1beta/openai"


def configure(get_session_id, broadcast, get_llm_chain, get_db_conn):
    """Call from app.py lifespan to inject shared state."""
    global _get_session_id, _broadcast, _get_llm_chain, _get_db_conn
    _get_session_id = get_session_id
    _broadcast = broadcast
    _get_llm_chain = get_llm_chain
    _get_db_conn = get_db_conn


# ─── Facilitator HTML views ───────────────────────────────────────────────────

@router.get("/facilitator", response_class=FileResponse)
async def facilitator_view():
    return FileResponse("facilitator.html")


@router.get("/corpus-admin", response_class=FileResponse)
async def corpus_admin_view():
    return FileResponse("corpus.html")


# ─── Live synthesis ──────────────────────────────────────────────────────────

@router.post("/v1/live-synthesis")
async def trigger_synthesis():
    result = await _run_synthesis()
    return JSONResponse(result)


async def _run_synthesis() -> dict:
    """Core synthesis. Called by route and background task."""
    session_id = _get_session_id()
    if not session_id:
        return {"error": "no_active_session", "themes": []}

    segments = await db_module.get_session_transcript(session_id)
    full_text = " ".join(s["text"] for s in segments if s.get("text"))
    transcript_excerpt = full_text[-16000:]
    if len(transcript_excerpt) < 100:
        return {"error": "transcript_too_short", "themes": []}

    # Corpus context (optional — skipped if embedding unavailable)
    corpus_passages: list[dict] = []
    conn = _get_db_conn()
    if conn is not None:
        loaded = await corpus_module.load_corpus(conn)
        if loaded:
            query_emb = await corpus_module.embed_text(transcript_excerpt[-3000:])
            if query_emb is not None:
                corpus_passages = corpus_module.search_corpus(query_emb, loaded, k=5)

    user_prompt = corpus_module.build_synthesis_user_prompt(transcript_excerpt, corpus_passages)
    chain = _get_llm_chain()
    synthesis = await _call_synthesis_llm(corpus_module.SYNTHESIS_SYSTEM_PROMPT, user_prompt, chain)

    message = {
        "type": "live_synthesis",
        "session_id": session_id,
        "generated_at": time.time(),
        **synthesis,
    }
    if _broadcast:
        await _broadcast(json.dumps(message))
    return synthesis


# ─── Q&A ─────────────────────────────────────────────────────────────────────

@router.post("/v1/qa")
async def qa(request: Request):
    body = await request.json()
    question = body.get("question", "").strip()
    if not question:
        return JSONResponse({"error": "No question"}, status_code=400)

    session_id = _get_session_id()
    if not session_id:
        return JSONResponse({"error": "No active session"}, status_code=400)

    segments = await db_module.get_session_transcript(session_id)
    transcript = " ".join(s["text"] for s in segments if s.get("text"))[-16000:]

    corpus_passages: list[dict] = []
    conn = _get_db_conn()
    if conn is not None:
        loaded = await corpus_module.load_corpus(conn)
        if loaded:
            emb = await corpus_module.embed_text(question)
            if emb is not None:
                corpus_passages = corpus_module.search_corpus(emb, loaded, k=3)

    user_prompt = corpus_module.build_qa_user_prompt(transcript, question, corpus_passages)
    chain = _get_llm_chain()
    answer = await _call_qa_llm(corpus_module.QA_SYSTEM_PROMPT, user_prompt, chain)
    return JSONResponse({"answer": answer})


# ─── Corpus CRUD ─────────────────────────────────────────────────────────────

@router.get("/v1/corpus")
async def list_corpus():
    conn = _get_db_conn()
    if conn is None:
        return JSONResponse({"docs": []})
    docs = await corpus_module.list_docs(conn)
    return JSONResponse({"docs": docs})


@router.patch("/v1/corpus/{doc_id}")
async def toggle_corpus_doc(doc_id: int, request: Request):
    body = await request.json()
    active = bool(body.get("active", True))
    conn = _get_db_conn()
    if conn is None:
        return JSONResponse({"error": "no db"}, status_code=500)
    await corpus_module.set_doc_active(conn, doc_id, active)
    return JSONResponse({"ok": True})


@router.delete("/v1/corpus/{doc_id}")
async def delete_corpus_doc(doc_id: int):
    conn = _get_db_conn()
    if conn is None:
        return JSONResponse({"error": "no db"}, status_code=500)
    await corpus_module.delete_doc(conn, doc_id)
    return JSONResponse({"ok": True})


# ─── PPTX export ─────────────────────────────────────────────────────────────

@router.post("/v1/sessions/{session_id}/export/synthesis-pptx")
async def export_synthesis_pptx(session_id: str):
    from pptx import Presentation

    segments = await db_module.get_session_transcript(session_id)
    full_text = " ".join(s["text"] for s in segments if s.get("text"))
    if not full_text:
        return JSONResponse({"error": "No transcript for session"}, status_code=404)

    user_prompt = corpus_module.build_synthesis_user_prompt(full_text[-16000:], [])
    chain = _get_llm_chain()
    synthesis = await _call_synthesis_llm(corpus_module.SYNTHESIS_SYSTEM_PROMPT, user_prompt, chain)
    themes = synthesis.get("themes", [])

    template_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets", "template_cap_pptx.pptx")
    prs = Presentation(template_path)

    # Cover slide — layout 21 "Title Subtitle"
    cover = prs.slides.add_slide(prs.slide_layouts[21])
    cover.placeholders[0].text = "Synthèse Workshop ASE"
    cover.placeholders[22].text = datetime.date.today().strftime("%d/%m/%Y")

    # One slide per theme — layout 35 "Content 3 Boxes"
    boxes_layout = prs.slide_layouts[35]
    for theme in themes:
        slide = prs.slides.add_slide(boxes_layout)
        slide.placeholders[0].text = theme.get("name", "")

        def fill_box(ph_idx: int, items: list[str], header: str):
            tf = slide.placeholders[ph_idx].text_frame
            tf.clear()
            tf.text = header
            for item in items:
                p = tf.add_paragraph()
                p.text = f"• {item}"

        fill_box(22, theme.get("alignment", []), "✓ Alignement")
        fill_box(35, theme.get("disagreement", []), "✗ Désaccord")
        fill_box(36, theme.get("unresolved", []), "? Non tranché")

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    date_str = datetime.date.today().strftime("%Y%m%d")
    filename = f"synthesis-{date_str}.pptx"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


# ─── LLM helpers ─────────────────────────────────────────────────────────────

async def _call_synthesis_llm(system: str, user: str, chain: list[dict]) -> dict:
    for tier in chain:
        try:
            raw = await _llm_call(tier, system, user)
            match = re.search(r'\{.*\}', raw, re.DOTALL)
            if match:
                return json.loads(match.group())
        except Exception:
            continue
    return {"themes": [], "error": "llm_unavailable"}


async def _call_qa_llm(system: str, user: str, chain: list[dict]) -> str:
    for tier in chain:
        try:
            return await _llm_call(tier, system, user)
        except Exception:
            continue
    return "LLM indisponible."


async def _llm_call(tier: dict, system: str, user: str) -> str:
    provider = tier["provider"]
    model = tier["model"]
    timeout = aiohttp.ClientTimeout(total=60)

    if provider == "anthropic":
        url = f"{ANTHROPIC_BASE_URL}/v1/messages"
        headers = {"anthropic-version": "2023-06-01", "content-type": "application/json"}
        if ANTHROPIC_AUTH_TOKEN:
            headers["Authorization"] = f"Bearer {ANTHROPIC_AUTH_TOKEN}"
        else:
            headers["x-api-key"] = ANTHROPIC_API_KEY
        body = {"model": model, "max_tokens": 2048, "system": system,
                "messages": [{"role": "user", "content": user}]}
        async with aiohttp.ClientSession() as s:
            async with s.post(url, headers=headers, json=body, timeout=timeout) as r:
                data = await r.json()
                return data["content"][0]["text"]

    # Hugin (Ollama) or Gemini — both use OpenAI-compatible /chat/completions
    if provider == "hugin":
        url = f"{HUGIN_BASE_URL}/v1/chat/completions"
        headers = {"Content-Type": "application/json"}
        if HUGIN_CF_ID and HUGIN_CF_SECRET:
            headers["CF-Access-Client-Id"] = HUGIN_CF_ID
            headers["CF-Access-Client-Secret"] = HUGIN_CF_SECRET
    else:  # gemini
        url = f"{GEMINI_BASE_URL}/chat/completions"
        headers = {"Authorization": f"Bearer {GEMINI_API_KEY}", "Content-Type": "application/json"}

    body = {
        "model": model, "max_tokens": 2048,
        "messages": [{"role": "system", "content": system}, {"role": "user", "content": user}],
    }
    async with aiohttp.ClientSession() as s:
        async with s.post(url, headers=headers, json=body, timeout=timeout) as r:
            data = await r.json()
            return data["choices"][0]["message"]["content"]
