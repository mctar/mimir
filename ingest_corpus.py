#!/usr/bin/env python3
"""
One-shot corpus ingestion for ASE workshop facilitator.
Usage: python ingest_corpus.py --dir /path/to/docs [--db livemind.db] [--chunk 2000] [--overlap 200] [--metadata corpus/metadata.yaml]

Supported: .txt .md .pdf .docx .pptx
Idempotent: skips already-ingested content (SHA-256 hash check); always syncs metadata.
"""

import argparse, asyncio, sys
from pathlib import Path

import aiosqlite
import corpus
from log import logger


def _parse_args():
    p = argparse.ArgumentParser(description="Ingest documents into ASE corpus.")
    p.add_argument("--dir", required=True, help="Directory of documents to ingest")
    p.add_argument("--db", default="livemind.db", help="Path to livemind.db")
    p.add_argument("--chunk", type=int, default=2000, help="Chunk size in chars")
    p.add_argument("--overlap", type=int, default=200, help="Overlap in chars")
    p.add_argument("--metadata", default="corpus/metadata.yaml",
                   help="YAML file with document-level metadata (label, role, key_messages, usages)")
    return p.parse_args()


def _load_metadata(path: str) -> dict:
    """Load corpus/metadata.yaml and return a dict keyed by filename."""
    try:
        import yaml
    except ImportError:
        logger.warning("PyYAML not installed — metadata will not be loaded. Run: pip install pyyaml")
        return {}
    p = Path(path)
    if not p.exists():
        logger.info(f"No metadata file at {path} — ingesting without document metadata.")
        return {}
    with open(p, encoding="utf-8") as f:
        data = yaml.safe_load(f) or {}
    result = {}
    for doc in data.get("documents", []):
        filename = doc.get("file", "")
        if filename:
            result[filename] = doc
    logger.info(f"Loaded metadata for {len(result)} documents from {path}")
    return result


def extract_text(path: Path) -> str:
    """Extract plain text from supported file types."""
    suffix = path.suffix.lower()
    if suffix in (".txt", ".md"):
        return path.read_text(encoding="utf-8", errors="replace")
    if suffix == ".pdf":
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            return "\n".join(page.extract_text() or "" for page in pdf.pages)
    if suffix == ".docx":
        from docx import Document
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs)
    if suffix == ".pptx":
        from pptx import Presentation
        prs = Presentation(path)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)
        return "\n".join(parts)
    raise ValueError(f"Unsupported file type: {suffix}")


def chunk_text(text: str, chunk_size: int, overlap: int) -> list[str]:
    """Split text into overlapping chunks."""
    if overlap >= chunk_size:
        raise ValueError(f"overlap ({overlap}) must be less than chunk_size ({chunk_size})")
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end].strip())
        start += chunk_size - overlap
    return [c for c in chunks if len(c) > 50]  # drop tiny tail chunks


async def ingest_file(db: aiosqlite.Connection, path: Path, chunk_size: int, overlap: int,
                      metadata_map: dict | None = None):
    print(f"  Processing {path.name}...", end="", flush=True)
    try:
        text = extract_text(path)
    except Exception as e:
        logger.error(f"Processing {path.name}: {e}")
        return

    meta = (metadata_map or {}).get(path.name, {})
    label = meta.get("label") or None
    role = meta.get("role") or None
    key_messages = meta.get("key_messages") or None
    usages = meta.get("usages") or None

    chunks = chunk_text(text, chunk_size, overlap)
    stored, updated, errored = 0, 0, 0
    for i, chunk in enumerate(chunks):
        title = f"{path.stem} [{i+1}/{len(chunks)}]"
        try:
            row_id = await corpus.store_doc(db, title, path.name, chunk,
                                            label=label, role=role,
                                            key_messages=key_messages, usages=usages)
            if row_id:
                stored += 1
            else:
                updated += 1  # chunk existed, metadata synced
        except Exception as e:
            logger.error(f"DB ERROR on chunk {i+1}: {e}")
            errored += 1
    summary = f"{stored} new, {updated} updated"
    if errored:
        summary += f", {errored} errors"
    logger.info(f"{path.name}: {summary}" + (f" [label: {label}]" if label else ""))


async def main():
    args = _parse_args()
    doc_dir = Path(args.dir)
    if not doc_dir.is_dir():
        logger.error(f"{doc_dir} is not a directory")
        sys.exit(1)

    supported = {".txt", ".md", ".pdf", ".docx", ".pptx"}
    files = [f for f in doc_dir.iterdir() if f.suffix.lower() in supported]
    if not files:
        logger.info(f"No supported files found in {doc_dir}")
        sys.exit(0)

    metadata_map = _load_metadata(args.metadata)

    logger.info(f"Found {len(files)} files in {doc_dir}")
    logger.info(f"DB: {args.db}")

    async with aiosqlite.connect(args.db) as db:
        await db.execute("""
            CREATE TABLE IF NOT EXISTS corpus_docs (
                id INTEGER PRIMARY KEY AUTOINCREMENT, title TEXT NOT NULL,
                source TEXT, content TEXT NOT NULL,
                content_hash TEXT NOT NULL UNIQUE,
                created_at REAL NOT NULL, active INTEGER NOT NULL DEFAULT 1
            )
        """)
        # Ensure metadata columns exist (standalone ingestion without app startup)
        for col_name, col_def in [
            ("label", "TEXT"), ("role", "TEXT"),
            ("key_messages", "TEXT"), ("usages", "TEXT"),
        ]:
            try:
                await db.execute(f"ALTER TABLE corpus_docs ADD COLUMN {col_name} {col_def}")
            except Exception:
                pass  # column already exists
        await db.commit()
        for f in sorted(files):
            await ingest_file(db, f, args.chunk, args.overlap, metadata_map=metadata_map)

    logger.info("Done.")


if __name__ == "__main__":
    asyncio.run(main())
